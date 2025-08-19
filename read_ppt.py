# read_ppt.py

from pptx import Presentation
from openai import OpenAI
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json
import re

from config import api_key, model

# OpenAIクライアント設定
client = OpenAI(api_key=api_key)


def is_image_placeholder(text):
    """画像のプレースホルダーかどうかを判定"""
    # 画像プレースホルダーとして扱うべきキーワード
    image_keywords = [
        "写真",
        "画像",
        "イメージ",
        "イラスト",
        "撮影",
        "全景写真",
    ]

    # 除外すべきパターン（画像プレースホルダーではない）
    exclude_patterns = [
        r"図っている",
        r"図る",
        r"図ら",
        r"図り",
        r"図れ",
        r"図ろう",
    ]

    # 除外パターンにマッチする場合は画像プレースホルダーではない
    for pattern in exclude_patterns:
        if re.search(pattern, text):
            return False

    # 画像キーワードにマッチする場合は画像プレースホルダー
    for keyword in image_keywords:
        if keyword in text:
            return True

    return False


def extract_text_from_shape(
    shape,
    slide_width=None,
    slide_height=None,
    elements=None,
    parent_left=0,
    parent_top=0,
):
    """シェイプからテキストを抽出（ページ外のテキストは除外）"""
    if elements is None:
        elements = []

    left = getattr(shape, "left", 0) + parent_left
    top = getattr(shape, "top", 0) + parent_top
    width = getattr(shape, "width", None)
    height = getattr(shape, "height", None)

    # シェイプ情報をデバッグファイルに出力
    shape_info = {
        "shape_type": str(shape.shape_type),
        "has_text_frame": shape.has_text_frame,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "text": getattr(shape, "text", "").strip() if hasattr(shape, "text") else "",
    }

    with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
        f.write(f"シェイプ情報: {shape_info}\n")

    # x軸方向の判定のみ
    if None not in (left, width, slide_width):
        # シェイプのx軸でページ内に存在する長さを計算
        x_in = min(left + width, slide_width) - max(left, 0)
        # シェイプの半分以上がページ内にあればページ内として判断
        if x_in <= 0 or x_in < width / 2:
            with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
                f.write(f"ページ外として除外: x_in={x_in}, width={width}\n")
            return elements

    # 画像の場合
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        elements.append({"text": "画像", "top": top, "left": left})
        return elements

    # テキストフレームを持つ場合
    elif shape.has_text_frame:
        text = shape.text.strip()
        if text:
            # 画像のプレースホルダーの場合は「画像」として扱う
            if is_image_placeholder(text):
                elements.append({"text": "画像", "top": top, "left": left})
                with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
                    f.write(f"画像プレースホルダーとして処理: {text}\n")
            else:
                elements.append({"text": text, "top": top, "left": left})
                with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
                    f.write(f"テキストとして処理: {text}\n")
        else:
            with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
                f.write("テキストフレームはあるが空文字列\n")

    # グループ化されたシェイプの場合
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            extract_text_from_shape(
                sub_shape, slide_width, slide_height, elements, left, top
            )

    # テーブルの場合
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    elements.append(
                        {"text": cell.text.strip(), "top": top, "left": left}
                    )

    # その他のシェイプタイプ
    else:
        with open("debug/shape_info.txt", "a", encoding="utf-8") as f:
            f.write(f"未対応のシェイプタイプ: {shape.shape_type}\n")

    return elements


# 除外するテキスト
def is_breadcrumb(text):
    return "HOME" in text and "＞" in text


# 除外するテキスト
def should_exclude_texts(texts, page_title):
    exclude_indices = set()
    if len(texts) >= 3:
        if (
            texts[0] == "ヘッダー"
            and texts[1] == page_title
            and is_breadcrumb(texts[2])
        ):
            exclude_indices.update([0, 1, 2])
    if texts and texts[-1] == "フッター":
        exclude_indices.add(len(texts) - 1)
    return exclude_indices


def extract_slide_texts(slide, slide_width, slide_height, page_title):
    elements = []
    for shape in slide.shapes:
        extract_text_from_shape(shape, slide_width, slide_height, elements)

    # 座標を正規化してソート（EMU座標をスライドサイズで割って正規化）
    def normalize_coordinates(element):
        # EMU座標をスライドサイズで正規化（0-1の範囲に）
        normalized_top = (
            element["top"] / slide_height if slide_height else element["top"]
        )
        normalized_left = (
            element["left"] / slide_width if slide_width else element["left"]
        )
        return (normalized_top, normalized_left)

    # 正規化された座標でソート
    elements.sort(key=normalize_coordinates)

    texts = [el["text"] for el in elements]
    exclude_indices = should_exclude_texts(texts, page_title)
    return [
        t
        for i, t in enumerate(texts)
        if i not in exclude_indices and not is_breadcrumb(t)
    ]


# rules.txtを読み込む
def load_rules():
    with open("rules/rules.txt", "r", encoding="utf-8") as f:
        return f.read()


# GPTに分類させる
def classify_texts_with_gpt(texts, rules_text):
    prompt = f"""{rules_text}

スライドの要素リスト：
{texts}

上記の内容をもとに、ルールに従ってHTML断片を生成してください。"""

    # デバッグ用：APIへの入力値をログに記録
    with open("debug/api_input.txt", "w", encoding="utf-8") as f:
        f.write("=== APIへの入力値 ===\n\n")
        f.write("【プロンプト】\n")
        f.write(prompt)
        f.write("\n\n【システムメッセージ】\n")
        f.write("あなたはパワーポイントスライドをHTML断片に変換するアシスタントです。")

    response = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "system",
                "content": "あなたはパワーポイントスライドをHTML断片に変換するアシスタントです。",
            },
            {
                "role": "user",
                "content": prompt,
            },
        ],
    )
    result = response.choices[0].message.content.strip()

    # デバッグ用：APIからの出力値をログに記録
    with open("debug/api_output.txt", "w", encoding="utf-8") as f:
        f.write("=== APIからの出力値 ===\n\n")
        f.write(result)

    return result


# パーツ名からルールファイルを取得
def load_parts_rule(parts_name):
    """
    parts_name（例: "カード"）から対応するルールファイル（例: rules/r_card.txt）の内容を返す
    見つからない場合は空文字列を返す
    """
    import json

    parts_list_path = os.path.join("rules", "parts_list.json")
    if not os.path.exists(parts_list_path):
        return ""
    with open(parts_list_path, "r", encoding="utf-8") as f:
        parts_list = json.load(f)
    for part in parts_list.get("parts", []):
        if part.get("parts_name") == parts_name:
            rule_file = os.path.join("rules", part.get("file_name"))
            if os.path.exists(rule_file):
                with open(rule_file, "r", encoding="utf-8") as rf:
                    return rf.read()
    return ""


# テキスト中にパーツ名が含まれていれば全て返す
def find_parts_name_in_texts(texts):
    import json

    parts_list_path = os.path.join("rules", "parts_list.json")
    if not os.path.exists(parts_list_path):
        return []
    with open(parts_list_path, "r", encoding="utf-8") as f:
        parts_list = json.load(f)
    parts_names = [part["parts_name"] for part in parts_list.get("parts", [])]
    found = set()
    for text in texts:
        text_clean = text.strip().lower()
        for name in parts_names:
            if name.lower() in text_clean:
                found.add(name)
    return list(found)


def generate_html_from_pptx(
    pptx_path, start_slide, end_slide, page_title, parts_name=None, output_filename=None
):
    start_slide = int(start_slide)
    end_slide = int(end_slide)
    os.makedirs("debug", exist_ok=True)

    # デバッグファイルを初期化
    with open("debug/shape_info.txt", "w", encoding="utf-8") as f:
        f.write("=== シェイプ情報デバッグ ===\n\n")

    prs = Presentation(pptx_path)
    common_rules = load_rules()
    slides = prs.slides
    start_slide_index = start_slide - 1
    end_slide_index = end_slide - 1
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 指定範囲の全スライド要素をまとめてリスト化
    all_texts = []
    for idx in range(start_slide_index, end_slide_index + 1):
        slide = slides[idx]
        texts = extract_slide_texts(slide, slide_width, slide_height, page_title)
        if texts:
            all_texts.extend(texts)

    # パーツ名自動判定
    if parts_name is None:
        parts_name = find_parts_name_in_texts(all_texts)

    # デバッグ情報をファイルに出力
    with open("debug/ppt_texts.json", "w", encoding="utf-8") as f:
        json.dump(all_texts, f, ensure_ascii=False, indent=2)

    # パーツルールを取得
    parts_rule = ""
    if isinstance(parts_name, list):
        for name in parts_name:
            rule = load_parts_rule(name)
            if rule:
                parts_rule += "\n" + rule
    elif parts_name:
        rule = load_parts_rule(parts_name)
        if rule:
            parts_rule = rule

    rules_text = common_rules
    if parts_rule:
        rules_text += "\n" + parts_rule

    # output_filenameをプロンプトに追加
    if output_filename:
        rules_text += f"\n\n【画像ファイル名の設定】\n- output_filename = {output_filename}\n- 画像ファイル名は「{{output_filename}}_img{{連番}}.jpg」の形式を使用してください。"

    # GPTに分類させる
    result = classify_texts_with_gpt(all_texts, rules_text)
    return result
